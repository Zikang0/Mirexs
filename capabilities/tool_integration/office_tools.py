"""
Office工具集成模块
提供Word、Excel、PowerPoint等Office文档的自动化操作
"""

import os
import logging
from typing import Dict, List, Any, Optional
from pathlib import Path
import pythoncom
from datetime import datetime

try:
    import win32com.client as win32
    from win32com.client import constants
    WIN32_AVAILABLE = True
except ImportError:
    WIN32_AVAILABLE = False

import openpyxl
from openpyxl import Workbook, load_workbook
from docx import Document
from pptx import Presentation
import pandas as pd

logger = logging.getLogger(__name__)

class WordProcessor:
    """Word文档处理器"""
    
    def __init__(self):
        self.doc = None
        self.app = None
    
    def create_document(self, template_path: Optional[str] = None) -> Dict[str, Any]:
        """创建新文档"""
        try:
            if template_path and os.path.exists(template_path):
                self.doc = Document(template_path)
            else:
                self.doc = Document()
            
            return {"success": True, "message": "文档创建成功"}
        except Exception as e:
            logger.error(f"文档创建失败: {e}")
            return {"success": False, "error": str(e)}
    
    def add_paragraph(self, text: str, style: str = 'Normal') -> Dict[str, Any]:
        """添加段落"""
        try:
            paragraph = self.doc.add_paragraph(text, style=style)
            return {"success": True, "message": "段落添加成功"}
        except Exception as e:
            logger.error(f"段落添加失败: {e}")
            return {"success": False, "error": str(e)}
    
    def add_heading(self, text: str, level: int = 1) -> Dict[str, Any]:
        """添加标题"""
        try:
            self.doc.add_heading(text, level=level)
            return {"success": True, "message": "标题添加成功"}
        except Exception as e:
            logger.error(f"标题添加失败: {e}")
            return {"success": False, "error": str(e)}
    
    def add_table(self, data: List[List[str]], rows: int, cols: int) -> Dict[str, Any]:
        """添加表格"""
        try:
            table = self.doc.add_table(rows=rows, cols=cols)
            
            for i, row_data in enumerate(data):
                for j, cell_data in enumerate(row_data):
                    if i < rows and j < cols:
                        table.cell(i, j).text = str(cell_data)
            
            return {"success": True, "message": "表格添加成功"}
        except Exception as e:
            logger.error(f"表格添加失败: {e}")
            return {"success": False, "error": str(e)}
    
    def save_document(self, filepath: str) -> Dict[str, Any]:
        """保存文档"""
        try:
            self.doc.save(filepath)
            return {"success": True, "message": f"文档保存成功: {filepath}"}
        except Exception as e:
            logger.error(f"文档保存失败: {e}")
            return {"success": False, "error": str(e)}
    
    def read_document(self, filepath: str) -> Dict[str, Any]:
        """读取文档内容"""
        try:
            if not os.path.exists(filepath):
                return {"success": False, "error": "文件不存在"}
            
            self.doc = Document(filepath)
            content = []
            
            for paragraph in self.doc.paragraphs:
                if paragraph.text.strip():
                    content.append(paragraph.text)
            
            # 提取表格内容
            tables = []
            for table in self.doc.tables:
                table_data = []
                for row in table.rows:
                    row_data = [cell.text for cell in row.cells]
                    table_data.append(row_data)
                tables.append(table_data)
            
            return {
                "success": True,
                "content": content,
                "tables": tables,
                "paragraph_count": len(content)
            }
        except Exception as e:
            logger.error(f"文档读取失败: {e}")
            return {"success": False, "error": str(e)}

class ExcelProcessor:
    """Excel处理器"""
    
    def __init__(self):
        self.workbook = None
        self.active_sheet = None
    
    def create_workbook(self) -> Dict[str, Any]:
        """创建工作簿"""
        try:
            self.workbook = Workbook()
            self.active_sheet = self.workbook.active
            return {"success": True, "message": "工作簿创建成功"}
        except Exception as e:
            logger.error(f"工作簿创建失败: {e}")
            return {"success": False, "error": str(e)}
    
    def load_workbook(self, filepath: str) -> Dict[str, Any]:
        """加载工作簿"""
        try:
            if not os.path.exists(filepath):
                return {"success": False, "error": "文件不存在"}
            
            self.workbook = load_workbook(filepath)
            self.active_sheet = self.workbook.active
            return {"success": True, "message": "工作簿加载成功"}
        except Exception as e:
            logger.error(f"工作簿加载失败: {e}")
            return {"success": False, "error": str(e)}
    
    def set_active_sheet(self, sheet_name: str) -> Dict[str, Any]:
        """设置活动工作表"""
        try:
            if sheet_name in self.workbook.sheetnames:
                self.active_sheet = self.workbook[sheet_name]
                return {"success": True, "message": f"活动工作表设置为: {sheet_name}"}
            else:
                return {"success": False, "error": f"工作表 '{sheet_name}' 不存在"}
        except Exception as e:
            logger.error(f"设置活动工作表失败: {e}")
            return {"success": False, "error": str(e)}
    
    def write_data(self, data: List[List[Any]], start_row: int = 1, start_col: int = 1) -> Dict[str, Any]:
        """写入数据"""
        try:
            for i, row in enumerate(data):
                for j, value in enumerate(row):
                    self.active_sheet.cell(
                        row=start_row + i,
                        column=start_col + j,
                        value=value
                    )
            
            return {"success": True, "message": "数据写入成功"}
        except Exception as e:
            logger.error(f"数据写入失败: {e}")
            return {"success": False, "error": str(e)}
    
    def read_range(self, start_row: int, start_col: int, end_row: int, end_col: int) -> Dict[str, Any]:
        """读取数据范围"""
        try:
            data = []
            for row in range(start_row, end_row + 1):
                row_data = []
                for col in range(start_col, end_col + 1):
                    cell_value = self.active_sheet.cell(row=row, column=col).value
                    row_data.append(cell_value)
                data.append(row_data)
            
            return {"success": True, "data": data}
        except Exception as e:
            logger.error(f"数据读取失败: {e}")
            return {"success": False, "error": str(e)}
    
    def add_chart(self, chart_type: str, data_range: str, title: str) -> Dict[str, Any]:
        """添加图表"""
        try:
            from openpyxl.chart import (
                BarChart, LineChart, PieChart, Reference
            )
            
            # 根据图表类型创建相应的图表对象
            if chart_type.lower() == 'bar':
                chart = BarChart()
            elif chart_type.lower() == 'line':
                chart = LineChart()
            elif chart_type.lower() == 'pie':
                chart = PieChart()
            else:
                return {"success": False, "error": f"不支持的图表类型: {chart_type}"}
            
            chart.title = title
            # 这里需要根据实际数据范围设置图表数据
            # 简化实现，实际使用时需要更复杂的逻辑
            
            self.active_sheet.add_chart(chart, "E5")
            return {"success": True, "message": "图表添加成功"}
        except Exception as e:
            logger.error(f"图表添加失败: {e}")
            return {"success": False, "error": str(e)}
    
    def save_workbook(self, filepath: str) -> Dict[str, Any]:
        """保存工作簿"""
        try:
            self.workbook.save(filepath)
            return {"success": True, "message": f"工作簿保存成功: {filepath}"}
        except Exception as e:
            logger.error(f"工作簿保存失败: {e}")
            return {"success": False, "error": str(e)}

class PowerPointProcessor:
    """PowerPoint处理器"""
    
    def __init__(self):
        self.presentation = None
    
    def create_presentation(self, template_path: Optional[str] = None) -> Dict[str, Any]:
        """创建演示文稿"""
        try:
            if template_path and os.path.exists(template_path):
                self.presentation = Presentation(template_path)
            else:
                self.presentation = Presentation()
            
            return {"success": True, "message": "演示文稿创建成功"}
        except Exception as e:
            logger.error(f"演示文稿创建失败: {e}")
            return {"success": False, "error": str(e)}
    
    def add_slide(self, layout_index: int = 0) -> Dict[str, Any]:
        """添加幻灯片"""
        try:
            slide_layout = self.presentation.slide_layouts[layout_index]
            slide = self.presentation.slides.add_slide(slide_layout)
            return {"success": True, "slide": slide}
        except Exception as e:
            logger.error(f"幻灯片添加失败: {e}")
            return {"success": False, "error": str(e)}
    
    def add_title_slide(self, title: str, subtitle: str = "") -> Dict[str, Any]:
        """添加标题幻灯片"""
        try:
            title_slide_layout = self.presentation.slide_layouts[0]
            slide = self.presentation.slides.add_slide(title_slide_layout)
            
            title_shape = slide.shapes.title
            subtitle_shape = slide.placeholders[1]
            
            title_shape.text = title
            subtitle_shape.text = subtitle
            
            return {"success": True, "message": "标题幻灯片添加成功"}
        except Exception as e:
            logger.error(f"标题幻灯片添加失败: {e}")
            return {"success": False, "error": str(e)}
    
    def add_content_slide(self, title: str, content: List[str]) -> Dict[str, Any]:
        """添加内容幻灯片"""
        try:
            content_layout = self.presentation.slide_layouts[1]
            slide = self.presentation.slides.add_slide(content_layout)
            
            title_shape = slide.shapes.title
            content_shape = slide.placeholders[1]
            
            title_shape.text = title
            
            # 添加内容
            text_frame = content_shape.text_frame
            text_frame.clear()
            
            for i, line in enumerate(content):
                if i == 0:
                    p = text_frame.paragraphs[0]
                else:
                    p = text_frame.add_paragraph()
                p.text = line
            
            return {"success": True, "message": "内容幻灯片添加成功"}
        except Exception as e:
            logger.error(f"内容幻灯片添加失败: {e}")
            return {"success": False, "error": str(e)}
    
    def save_presentation(self, filepath: str) -> Dict[str, Any]:
        """保存演示文稿"""
        try:
            self.presentation.save(filepath)
            return {"success": True, "message": f"演示文稿保存成功: {filepath}"}
        except Exception as e:
            logger.error(f"演示文稿保存失败: {e}")
            return {"success": False, "error": str(e)}

class OfficeAutomation:
    """Office自动化管理器"""
    
    def __init__(self):
        self.word_processor = WordProcessor()
        self.excel_processor = ExcelProcessor()
        self.powerpoint_processor = PowerPointProcessor()
    
    def convert_document(self, source_path: str, target_format: str) -> Dict[str, Any]:
        """转换文档格式"""
        try:
            if not os.path.exists(source_path):
                return {"success": False, "error": "源文件不存在"}
            
            base_name = os.path.splitext(source_path)[0]
            target_path = f"{base_name}.{target_format}"
            
            # 根据文件类型执行转换
            if source_path.endswith('.docx'):
                if target_format == 'pdf':
                    # 这里需要安装额外的库或使用COM接口
                    return {"success": False, "error": "PDF转换功能需要额外配置"}
            
            return {"success": True, "target_path": target_path}
        except Exception as e:
            logger.error(f"文档转换失败: {e}")
            return {"success": False, "error": str(e)}
    
    def batch_process_documents(self, folder_path: str, operation: str, **kwargs) -> Dict[str, Any]:
        """批量处理文档"""
        try:
            if not os.path.exists(folder_path):
                return {"success": False, "error": "文件夹不存在"}
            
            results = []
            for filename in os.listdir(folder_path):
                file_path = os.path.join(folder_path, filename)
                
                if operation == "convert":
                    result = self.convert_document(file_path, kwargs.get('target_format', 'pdf'))
                elif operation == "extract_text":
                    if filename.endswith('.docx'):
                        result = self.word_processor.read_document(file_path)
                    else:
                        result = {"success": False, "error": "不支持的格式"}
                
                results.append({
                    "filename": filename,
                    "result": result
                })
            
            return {"success": True, "results": results}
        except Exception as e:
            logger.error(f"批量处理失败: {e}")
            return {"success": False, "error": str(e)}

